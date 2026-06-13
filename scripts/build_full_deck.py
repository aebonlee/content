#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
통합 강의 덱(.pptx) 생성기 — 7개 모듈을 하나의 캠페인 내러티브로 연결.
러닝 케이스(Night Off)를 하루 동안 끌고 가며, 모듈마다 실제 입력→출력 예시와
좋은예 vs 나쁜예 비교, 개념 심화 슬라이드를 더해 분량과 깊이를 키운다.

실행: python3 scripts/build_full_deck.py
출력: docs/slides/pptx/AI홍보실무_통합강의.pptx
"""
import os
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from build_slides import (Deck, bg, rect, textbox, para, set_font,
                          PAPER, INK, PIGMENT, ACCENT, WASH, INK70, WHITE, PAPER_D, CODEBG,
                          TITLE_FONT, BODY_FONT, MONO_FONT)
from slides_data import MODULES
from slides_enrich import ENRICH, CONNECT, JOURNEY, CASE

# 비교 슬라이드 색(교육용 대비 — 절제된 톤)
BAD_BG, BAD_TX = RGBColor(0xF7, 0xEA, 0xEA), RGBColor(0xA6, 0x33, 0x33)
GOOD_BG, GOOD_TX = RGBColor(0xE7, 0xF1, 0xE9), RGBColor(0x2C, 0x6B, 0x44)


class FullDeck(Deck):
    def __init__(self):
        super().__init__(mod=None)

    # ---- 표지 / 여정 / 모듈 구분 / 마무리 --------------------------------
    def course_cover(self):
        s = self.blank(INK)
        for i in range(1, 8):
            rect(s, Inches(0.0), Inches(i * 0.95), Inches(13.333), Pt(0.5), RGBColor(0x22, 0x30, 0x50))
        tf = textbox(s, Inches(0.95), Inches(1.7), Inches(11.5), Inches(4.2))
        para(tf, "ONE-DAY WORKSHOP · AI MARKETING · 8 HOURS", 14, ACCENT, MONO_FONT, first=True, space_after=14)
        para(tf, "AI로 만드는 홍보 실무", 46, WHITE, TITLE_FONT, bold=True, space_after=6, line=1.05)
        para(tf, "콘텐츠 · 이미지 · 페이지까지 자동화 완성", 22, RGBColor(0xC7, 0xD2, 0xE6), TITLE_FONT, space_after=18)
        para(tf, "7개 모듈 · 14세션 · 하루 동안 하나의 캠페인을 처음부터 끝까지", 15, RGBColor(0x9F, 0xAE, 0xC8), BODY_FONT)
        box = rect(s, Inches(0.95), Inches(6.05), Inches(11.45), Inches(0.9), RGBColor(0x1B, 0x2A, 0x49))
        tfx = box.text_frame; tfx.word_wrap = True
        tfx.margin_left = Inches(0.3); tfx.vertical_anchor = MSO_ANCHOR.MIDDLE
        para(tfx, [("오늘의 러닝 케이스   ", {'color': ACCENT, 'bold': True, 'font': MONO_FONT, 'size': 13}),
                   (CASE, {'color': RGBColor(0xD7, 0xE0, 0xF0), 'size': 14})], 14, first=True)

    def journey(self):
        s = self.blank()
        self.kicker_title(s, "TODAY'S JOURNEY", "오늘의 여정 — 7개 모듈, 하나의 흐름")
        tf = textbox(s, Inches(1.0), Inches(1.95), Inches(11.3), Inches(5.0))
        for i, (tag, q, desc) in enumerate(JOURNEY):
            para(tf, [(tag + "  ", {'color': ACCENT, 'bold': True, 'font': MONO_FONT, 'size': 16}),
                      (q + "  ", {'color': INK, 'bold': True, 'size': 16}),
                      ("— " + desc, {'color': INK70, 'size': 14})],
                 15, first=(i == 0), space_before=7, space_after=6, line=1.2)
        box = rect(s, Inches(1.0), Inches(6.45), Inches(11.3), Inches(0.55), WASH)
        tfx = box.text_frame; tfx.vertical_anchor = MSO_ANCHOR.MIDDLE; tfx.margin_left = Inches(0.3)
        para(tfx, "각 모듈은 앞 모듈의 결과물을 입력으로 받아 이어집니다 — 따로가 아니라 하나의 캠페인.",
             13, PIGMENT, MONO_FONT, first=True)

    def module_divider(self, mod):
        s = self.blank(PIGMENT)
        rect(s, Inches(0.9), Inches(1.7), Inches(1.1), Inches(0.14), ACCENT)
        tf = textbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.6))
        para(tf, "MODULE %d · %s · %s" % (mod['no'], mod['period'], mod['time']),
             14, RGBColor(0xBF, 0xD2, 0xEC), MONO_FONT, first=True, space_after=10)
        para(tf, "M%d. %s" % (mod['no'], mod['title']), 38, WHITE, TITLE_FONT, bold=True, space_after=12, line=1.1)
        para(tf, mod['goal'], 15, RGBColor(0xDB, 0xE6, 0xF6), BODY_FONT, line=1.4)
        prev, this, nxt = CONNECT[mod['no']]
        box = rect(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(1.15), RGBColor(0x18, 0x40, 0x70))
        tfx = box.text_frame; tfx.word_wrap = True
        tfx.margin_left = Inches(0.35); tfx.margin_top = Inches(0.18); tfx.vertical_anchor = MSO_ANCHOR.MIDDLE
        para(tfx, [("이전  ", {'color': RGBColor(0x9F, 0xC0, 0xE6), 'font': MONO_FONT, 'size': 12, 'bold': True}),
                   (prev, {'color': WHITE, 'size': 13})], 13, first=True, space_after=3)
        para(tfx, [("이번  ", {'color': ACCENT, 'font': MONO_FONT, 'size': 12, 'bold': True}),
                   (this, {'color': WHITE, 'size': 13, 'bold': True})], 13, space_after=3)
        para(tfx, [("다음  ", {'color': RGBColor(0x9F, 0xC0, 0xE6), 'font': MONO_FONT, 'size': 12, 'bold': True}),
                   (nxt, {'color': WHITE, 'size': 13})], 13)

    def concept_examples(self, concept_ex):
        s = self.blank()
        self.kicker_title(s, "개념 깊이 보기", "개념을 예시로 — 우리 케이스에 적용하면")
        tf = textbox(s, Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.9))
        items = list(concept_ex.items())
        for i, (term, ex) in enumerate(items):
            para(tf, [("● ", {'color': PIGMENT, 'bold': True, 'size': 15}),
                      (term, {'color': INK, 'bold': True, 'size': 16})],
                 16, first=(i == 0), space_before=10, space_after=2, line=1.2)
            para(tf, [("    예) ", {'color': ACCENT, 'bold': True, 'font': MONO_FONT, 'size': 13}),
                      (ex, {'color': INK70, 'size': 14})],
                 14, space_before=0, space_after=4, line=1.25)

    def deep_slide(self, title, body):
        s = self.blank()
        self.kicker_title(s, "개념 심화", title)
        box = rect(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.0), WASH)
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.5); tf.margin_right = Inches(0.5)
        tf.margin_top = Inches(0.4); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para(tf, body, 18, INK, BODY_FONT, first=True, line=1.5)

    def io_slide(self, io):
        s = self.blank()
        self.kicker_title(s, "실제 입력 → 출력", io['title'])
        # 입력(프롬프트) — 짙은 박스
        ib = rect(s, Inches(1.0), Inches(1.95), Inches(11.3), Inches(2.05), CODEBG)
        rect(s, Inches(1.0), Inches(1.95), Inches(2.2), Inches(0.36), RGBColor(0x24, 0x3A, 0x60))
        lt = ib.text_frame; lt.word_wrap = True; lt.vertical_anchor = MSO_ANCHOR.TOP
        lt.margin_left = Inches(0.4); lt.margin_right = Inches(0.35); lt.margin_top = Inches(0.55)
        for i, ln in enumerate(io['prompt'].split('\n')):
            p = lt.paragraphs[0] if i == 0 else lt.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(0); p.space_after = Pt(1); p.line_spacing = 1.12
            r = p.add_run(); r.text = ln or ' '; r.font.size = Pt(12.5)
            r.font.color.rgb = RGBColor(0xE8, 0xEE, 0xFA); set_font(r, MONO_FONT)
        tlab = textbox(s, Inches(1.2), Inches(2.0), Inches(3), Inches(0.3))
        para(tlab, "▸ 입력 프롬프트", 12, ACCENT, MONO_FONT, first=True)
        # 출력 — 밝은 박스
        ob = rect(s, Inches(1.0), Inches(4.15), Inches(11.3), Inches(2.75), RGBColor(0xFB, 0xFC, 0xFE))
        ob.line.color.rgb = PAPER_D; ob.line.width = Pt(1)
        rect(s, Inches(1.0), Inches(4.15), Inches(2.0), Inches(0.36), WASH)
        ot = ob.text_frame; ot.word_wrap = True; ot.vertical_anchor = MSO_ANCHOR.TOP
        ot.margin_left = Inches(0.4); ot.margin_right = Inches(0.35); ot.margin_top = Inches(0.55)
        for i, ln in enumerate(io['output'].split('\n')):
            p = ot.paragraphs[0] if i == 0 else ot.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(0); p.space_after = Pt(1); p.line_spacing = 1.15
            mono = any(c in ln for c in '┌┐└┘├┤┬┴┼─│')
            r = p.add_run(); r.text = ln or ' '
            r.font.size = Pt(12 if mono else 13.5)
            r.font.color.rgb = INK if not ln.strip().startswith('→') else PIGMENT
            set_font(r, MONO_FONT if mono else BODY_FONT)
        olab = textbox(s, Inches(1.2), Inches(4.2), Inches(3), Inches(0.3))
        para(olab, "▸ 나온 결과(예)", 12, PIGMENT, MONO_FONT, first=True)

    def compare_slide(self, c):
        s = self.blank()
        self.kicker_title(s, "좋은예 vs 나쁜예", c['title'])
        # 왼쪽 나쁜예
        lb = rect(s, Inches(1.0), Inches(2.1), Inches(5.5), Inches(4.5), BAD_BG)
        lt = lb.text_frame; lt.word_wrap = True; lt.vertical_anchor = MSO_ANCHOR.TOP
        lt.margin_left = Inches(0.4); lt.margin_right = Inches(0.35); lt.margin_top = Inches(0.35)
        para(lt, c['bad_label'], 17, BAD_TX, TITLE_FONT, bold=True, first=True, space_after=10)
        for ln in c['bad']:
            para(lt, [("· ", {'color': BAD_TX, 'bold': True}), (ln, {'color': INK, 'size': 14})],
                 14, space_before=4, space_after=4, line=1.3)
        # 오른쪽 좋은예
        rb = rect(s, Inches(6.83), Inches(2.1), Inches(5.5), Inches(4.5), GOOD_BG)
        rt = rb.text_frame; rt.word_wrap = True; rt.vertical_anchor = MSO_ANCHOR.TOP
        rt.margin_left = Inches(0.4); rt.margin_right = Inches(0.35); rt.margin_top = Inches(0.35)
        para(rt, c['good_label'], 17, GOOD_TX, TITLE_FONT, bold=True, first=True, space_after=10)
        for ln in c['good']:
            para(rt, [("· ", {'color': GOOD_TX, 'bold': True}), (ln, {'color': INK, 'size': 14})],
                 14, space_before=4, space_after=4, line=1.3)

    def closing(self):
        s = self.blank(INK)
        rect(s, Inches(0.9), Inches(1.5), Inches(1.1), Inches(0.14), ACCENT)
        tf = textbox(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(4.8))
        para(tf, "WRAP-UP", 14, ACCENT, MONO_FONT, first=True, space_after=10)
        para(tf, "하루의 결과물이, 내일의 자산으로", 34, WHITE, TITLE_FONT, bold=True, space_after=16, line=1.1)
        para(tf, "오늘 'Night Off' 캠페인을 처음부터 끝까지 만들었습니다 —", 16, RGBColor(0xC7, 0xD2, 0xE6), BODY_FONT, space_after=8, line=1.4)
        for t in ["메시지 하우스(M2) → 콘텐츠 세트(M3) → 시각자료(M4)",
                  "→ UI 시안(M5) → 홍보 페이지(M6) → 자동화 루틴(M7)"]:
            para(tf, t, 16, WHITE, BODY_FONT, bold=True, space_before=2, space_after=2, line=1.4)
        para(tf, "이제 제품명만 바꾸면, 다음 캠페인은 처음부터 시작하지 않습니다.",
             15, RGBColor(0x9F, 0xAE, 0xC8), BODY_FONT, space_before=16)


def build_module(d, mod):
    d.module_divider(mod)
    d.bullets("학습 목표", "이 모듈에서 할 수 있게 되는 것", mod['objectives'])
    for idx, ses in enumerate(mod['sessions'], 1):
        e = ENRICH.get((mod['no'], idx), {})
        d.section("세션 %d" % idx, ses['title'], ses['lead'])
        d.bullets("핵심 개념", "핵심 개념", ses['concepts'])
        if e.get('concept_ex'):
            d.concept_examples(e['concept_ex'])
        for t, b in e.get('deep', []):
            d.deep_slide(t, b)
        d.bullets("실습 절차", "실습 절차 — 직접 해봅니다", ses['steps'], numbered=True)
        d.prompt(ses['prompt'][0], ses['prompt'][1])
        if e.get('io'):
            d.io_slide(e['io'])
        if e.get('compare'):
            d.compare_slide(e['compare'])
        d.callout("정리", ses['title'], pitfalls=ses.get('pitfalls'), output=ses.get('output'))
    d.bullets("도구 가이드", "이 모듈에서 쓰는 도구", mod['tools'])
    d.bullets("CHECKLIST", "모듈 체크리스트", mod['checklist'], bullet='☐')
    d.mission(mod['mission'][0], mod['mission'][1], mod['mission'][2], mod.get('next'))


def main():
    d = FullDeck()
    d.course_cover()
    d.journey()
    for mod in MODULES:
        build_module(d, mod)
    d.closing()
    here = os.path.dirname(os.path.abspath(__file__))
    outdir = os.path.join(here, '..', 'docs', 'slides', 'pptx')
    os.makedirs(outdir, exist_ok=True)
    path = os.path.join(outdir, 'AI홍보실무_통합강의.pptx')
    d.save(path)
    n = len(d.prs.slides._sldIdLst)
    print("✓ %s — 총 %d 슬라이드" % (os.path.basename(path), n))


if __name__ == '__main__':
    main()

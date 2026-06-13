#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
강의안 → PowerPoint(.pptx) 생성기.
src/config/lectures.js + site.js의 강의안을 옮긴 데이터로, 사이트와 동일한
아트디렉션(다크 남색·학술 블루·형광 호박 / Serif 제목)을 입힌 16:9 덱을
모듈별(M1~M7)로 docs/slides/pptx/ 에 생성한다.

실행:  python3 scripts/build_slides.py
의존:  python-pptx  (pip install --user python-pptx)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- 팔레트 (CLAUDE.md 아트디렉션) -----------------------------------------
PAPER   = RGBColor(0xF4, 0xF6, 0xFB)
INK     = RGBColor(0x14, 0x21, 0x3D)
PIGMENT = RGBColor(0x1D, 0x4E, 0x89)
ACCENT  = RGBColor(0xE8, 0xA3, 0x3D)
WASH    = RGBColor(0xDC, 0xE6, 0xF2)
INK70   = RGBColor(0x57, 0x61, 0x76)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
PAPER_D = RGBColor(0xE9, 0xEC, 0xF4)
CODEBG  = RGBColor(0x10, 0x1A, 0x30)  # 프롬프트 박스(짙은 남색)

TITLE_FONT = 'Noto Serif KR'
BODY_FONT  = 'Pretendard'
MONO_FONT  = 'JetBrains Mono'

EMU_W, EMU_H = Inches(13.333), Inches(7.5)


def set_font(run, name):
    """라틴·한글(ea)·기타(cs) 글꼴을 함께 지정해 한글도 의도한 폰트로."""
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set('typeface', name)


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    return tf


def para(tf, runs, size, color=INK, font=BODY_FONT, bold=False, first=False,
         space_before=6, space_after=4, align=PP_ALIGN.LEFT, line=1.18):
    """runs: 문자열 또는 [(text, {opts})] 리스트."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    try:
        p.line_spacing = line
    except Exception:
        pass
    if isinstance(runs, str):
        runs = [(runs, {})]
    for text, o in runs:
        r = p.add_run(); r.text = text
        r.font.size = Pt(o.get('size', size))
        r.font.bold = o.get('bold', bold)
        r.font.color.rgb = o.get('color', color)
        set_font(r, o.get('font', font))
    return p


# =========================================================================
# 슬라이드 빌더
# =========================================================================
class Deck:
    def __init__(self, mod):
        self.prs = Presentation()
        self.prs.slide_width = EMU_W
        self.prs.slide_height = EMU_H
        self.mod = mod

    def blank(self, color=PAPER):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg(s, color)
        return s

    def kicker_title(self, s, kicker, title, title_color=INK, top=Inches(0.55)):
        rect(s, Inches(0.7), top + Inches(0.02), Inches(0.18), Inches(0.42), ACCENT)
        tf = textbox(s, Inches(1.0), top, Inches(11.6), Inches(1.1))
        para(tf, kicker.upper(), 13, PIGMENT, MONO_FONT, first=True, space_before=0, space_after=2)
        para(tf, title, 30, title_color, TITLE_FONT, bold=True, space_before=0)

    # --- 표지 -----------------------------------------------------------
    def title_slide(self):
        m = self.mod
        s = self.blank(INK)
        # 모눈 느낌의 옅은 라인
        for i in range(1, 8):
            rect(s, Inches(0.0), Inches(i * 0.95), EMU_W, Pt(0.5), RGBColor(0x22, 0x30, 0x50))
        tf = textbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(3.5))
        para(tf, m['period'].upper() + ' · ' + m['time'] + ('  ·  2시간' if m['hours'] == 2 else '  ·  1시간'),
             14, ACCENT, MONO_FONT, first=True, space_after=10)
        para(tf, "M%d. %s" % (m['no'], m['title']), 40, WHITE, TITLE_FONT, bold=True, space_after=14, line=1.1)
        para(tf, m['goal'], 16, RGBColor(0xC7, 0xD2, 0xE6), BODY_FONT, line=1.4)
        tf2 = textbox(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.6))
        para(tf2, "AI로 만드는 홍보 실무 — 1일 워크숍 학습자료", 12, RGBColor(0x8A, 0x99, 0xB5),
             MONO_FONT, first=True)

    # --- 목표 / 목차 등 불릿 -------------------------------------------
    def bullets(self, kicker, title, items, numbered=False, color=INK, bullet='•'):
        s = self.blank()
        self.kicker_title(s, kicker, title)
        tf = textbox(s, Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.9))
        for i, it in enumerate(items):
            mark = ("%d  " % (i + 1)) if numbered else (bullet + "  ")
            mc = ACCENT if numbered else PIGMENT
            if isinstance(it, tuple):  # (term, desc)
                para(tf, [(mark, {'color': mc, 'bold': True, 'font': MONO_FONT, 'size': 17}),
                          (it[0] + " — ", {'bold': True, 'color': INK, 'size': 17}),
                          (it[1], {'color': INK70, 'size': 15})],
                     16, first=(i == 0), space_before=8, space_after=6, line=1.25)
            else:
                para(tf, [(mark, {'color': mc, 'bold': True, 'font': MONO_FONT, 'size': 17}),
                          (it, {'color': color, 'size': 17})],
                     17, first=(i == 0), space_before=8, space_after=6, line=1.25)
        return s

    # --- 세션 구분 ------------------------------------------------------
    def section(self, kicker, title, lead):
        s = self.blank(PIGMENT)
        rect(s, Inches(0.9), Inches(2.7), Inches(0.9), Inches(0.12), ACCENT)
        tf = textbox(s, Inches(0.9), Inches(3.0), Inches(11.5), Inches(2.6))
        para(tf, kicker.upper(), 14, RGBColor(0xBF, 0xD2, 0xEC), MONO_FONT, first=True, space_after=8)
        para(tf, title, 32, WHITE, TITLE_FONT, bold=True, space_after=12, line=1.12)
        para(tf, lead, 16, RGBColor(0xDB, 0xE6, 0xF6), BODY_FONT, line=1.4)

    # --- 프롬프트 박스 --------------------------------------------------
    def prompt(self, title, code):
        s = self.blank()
        self.kicker_title(s, "실습 프롬프트", title)
        box = rect(s, Inches(1.0), Inches(2.05), Inches(11.3), Inches(4.85), CODEBG)
        # 상단 라벨 점
        rect(s, Inches(1.25), Inches(2.3), Inches(0.12), Inches(0.12), ACCENT)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.45); tf.margin_right = Inches(0.4)
        tf.margin_top = Inches(0.55); tf.margin_bottom = Inches(0.35)
        for i, ln in enumerate(code.split('\n')):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(0); p.space_after = Pt(2); p.line_spacing = 1.15
            r = p.add_run(); r.text = ln if ln else ' '
            r.font.size = Pt(14)
            r.font.color.rgb = RGBColor(0xE8, 0xEE, 0xFA)
            set_font(r, MONO_FONT)

    # --- 결과물/실수 콜아웃 묶음 ---------------------------------------
    def callout(self, kicker, title, pitfalls=None, output=None):
        s = self.blank()
        self.kicker_title(s, kicker, title)
        top = 2.1
        if pitfalls:
            box = rect(s, Inches(1.0), Inches(top), Inches(11.3), Inches(0.05 + 0.62 * len(pitfalls) + 0.5),
                       RGBColor(0xFB, 0xF1, 0xDE))
            tf = box.text_frame; tf.word_wrap = True
            tf.margin_left = Inches(0.4); tf.margin_right = Inches(0.4)
            tf.margin_top = Inches(0.25); tf.margin_bottom = Inches(0.2)
            para(tf, [("⚠  흔한 실수", {'bold': True, 'color': RGBColor(0xB5, 0x7A, 0x1A), 'size': 15, 'font': MONO_FONT})],
                 15, first=True, space_after=6)
            for p in pitfalls:
                para(tf, [("·  ", {'color': ACCENT, 'bold': True}), (p, {'color': INK, 'size': 14})],
                     14, space_before=2, space_after=2, line=1.2)
            top += 0.05 + 0.62 * len(pitfalls) + 0.7
        if output:
            box = rect(s, Inches(1.0), Inches(top), Inches(11.3), Inches(1.1), WASH)
            tf = box.text_frame; tf.word_wrap = True
            tf.margin_left = Inches(0.4); tf.margin_right = Inches(0.4)
            tf.margin_top = Inches(0.22); tf.margin_bottom = Inches(0.2)
            para(tf, [("📦  결과물   ", {'bold': True, 'color': PIGMENT, 'size': 15, 'font': MONO_FONT}),
                      (output, {'color': INK, 'size': 16, 'bold': True})],
                 16, first=True, line=1.25)

    # --- 미션 마무리 ----------------------------------------------------
    def mission(self, mno, mtitle, tip, nxt):
        s = self.blank(INK)
        rect(s, Inches(0.9), Inches(1.5), Inches(0.9), Inches(0.12), ACCENT)
        tf = textbox(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(4.6))
        para(tf, "MISSION · 도장깨기", 14, ACCENT, MONO_FONT, first=True, space_after=10)
        para(tf, "미션 %02d — %s" % (mno, mtitle), 30, WHITE, TITLE_FONT, bold=True, space_after=12, line=1.1)
        para(tf, "💡 " + tip, 16, RGBColor(0xC7, 0xD2, 0xE6), BODY_FONT, space_after=24, line=1.4)
        if nxt:
            para(tf, [("▶  다음  ", {'color': ACCENT, 'bold': True, 'font': MONO_FONT, 'size': 15}),
                      (nxt, {'color': WHITE, 'size': 18, 'bold': True})], 18, line=1.3)

    def save(self, path):
        self.prs.save(path)


# =========================================================================
# 모듈 빌드 흐름
# =========================================================================
def build(mod, outdir):
    d = Deck(mod)
    d.title_slide()
    d.bullets("학습 목표", "학습 목표", mod['objectives'])
    d.bullets("AGENDA", "오늘의 흐름", [s['title'] for s in mod['sessions']], numbered=True)
    for idx, ses in enumerate(mod['sessions'], 1):
        d.section("세션 %d" % idx, ses['title'], ses['lead'])
        d.bullets("핵심 개념", "핵심 개념", ses['concepts'])
        d.bullets("실습 절차", "실습 절차", ses['steps'], numbered=True)
        d.prompt(ses['prompt'][0], ses['prompt'][1])
        d.callout("정리", ses['title'], pitfalls=ses.get('pitfalls'), output=ses.get('output'))
    d.bullets("도구 가이드", "도구 가이드", mod['tools'])
    d.bullets("CHECKLIST", "모듈 체크리스트", mod['checklist'], bullet='☐')
    d.mission(mod['mission'][0], mod['mission'][1], mod['mission'][2], mod.get('next'))
    fn = "M%d_%s.pptx" % (mod['no'], mod['slug'])
    path = os.path.join(outdir, fn)
    d.save(path)
    n = len(d.prs.slides._sldIdLst)
    return fn, n


def main():
    from slides_data import MODULES
    here = os.path.dirname(os.path.abspath(__file__))
    outdir = os.path.join(here, '..', 'docs', 'slides', 'pptx')
    os.makedirs(outdir, exist_ok=True)
    print("생성 위치:", os.path.normpath(outdir))
    for mod in MODULES:
        fn, n = build(mod, outdir)
        print("  ✓ %-26s %2d 슬라이드" % (fn, n))


if __name__ == '__main__':
    main()
